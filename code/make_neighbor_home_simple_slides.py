from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


ROOT = Path("/Users/jaimebenavides/claude_cowork/green_spaces_community_severance")
OUT = ROOT / "output"
TEMPLATE = OUT / "report" / "SPH_Powerpoint-Template-2026.pptx"
PPTX_PATH = OUT / "neighbor_home_indicator_and_models_2019_full_year_brown_sph.pptx"

IMG_PRIMARY_LA = OUT / "models_result_neighbor_visit_primary_fit_city_la_2019_full_year.png"
IMG_PRIMARY_NYC = OUT / "models_result_neighbor_visit_primary_fit_city_nyc_2019_full_year.png"
IMG_SHARE_LA = OUT / "models_result_neighbor_visit_share_fit_city_la_2019_full_year.png"
IMG_SHARE_NYC = OUT / "models_result_neighbor_visit_share_fit_city_nyc_2019_full_year.png"


def delete_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def clear_template_slides(prs):
    for idx in range(len(prs.slides) - 1, -1, -1):
        delete_slide(prs, idx)


def get_layout(prs, name, fallback_idx):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[fallback_idx]


def find_placeholder(slide, idx):
    try:
        return slide.placeholders[idx]
    except Exception:
        return None


def set_text(shape, text, size=20, bold=False, color=(30, 30, 30), align=PP_ALIGN.LEFT):
    if shape is None:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)


def add_bullets_to_placeholder(shape, bullets, size=20):
    if shape is None:
        return
    tf = shape.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(8)


def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(30, 30, 30)
    return box


def add_image(slide, img_path, left, top, width):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(width))


prs = Presentation(str(TEMPLATE))
clear_template_slides(prs)

layout_title = get_layout(prs, "TITLE", 0)
layout_content = get_layout(prs, "Title + Content", 2)
layout_object = get_layout(prs, "OBJECT", 1)


# Slide 1: title
slide = prs.slides.add_slide(layout_title)
set_text(find_placeholder(slide, 0), "Neighbor-home indicator and model results", size=28, bold=True)
set_text(
    find_placeholder(slide, 1),
    "Green spaces and community severance study\nFull-year 2019, NYC and LA",
    size=18
)


# Slide 2: construction
slide = prs.slides.add_slide(layout_content)
set_text(find_placeholder(slide, 0), "How the indicator is built", size=24, bold=True)
add_bullets_to_placeholder(
    find_placeholder(slide, 1),
    [
        "Start with one destination CBG-month from Dewey.",
        "Build a 0.5-mile buffer around the destination polygon.",
        "Select neighboring CBG polygons that intersect that buffer.",
        "Sum device counts from home CBGs inside that neighboring set.",
        "Compute count and share: NH_i and S_i = NH_i / HOME_DEVICE_COUNTS_TOTAL_PARSED."
    ],
    size=20
)
add_textbox(
    slide,
    "Important: the buffer is polygon-based, not centroid-based.",
    0.9, 6.3, 10.8, 0.4, font_size=15
)


# Slide 3: translation to study unit
slide = prs.slides.add_slide(layout_content)
set_text(find_placeholder(slide, 0), "How it enters this study", size=24, bold=True)
add_bullets_to_placeholder(
    find_placeholder(slide, 1),
    [
        "Destination CBGs were assigned to tracts using the first 11 digits of AREA.",
        "CBG-month values were summed to tract-month values within each city.",
        "Tract-month values were averaged across all 12 months of 2019.",
        "Primary outcome: neighbor_visit_count_annual_avg.",
        "Sensitivity outcome: neighbor_visit_share_annual_avg."
    ],
    size=20
)
add_textbox(
    slide,
    "Main model: negative binomial GAM with offset(log(HOME_DEVICE_COUNTS_TOTAL_PARSED))\nAdjusted covariates: perc.black, perc.hisp, perc.pov, pop_dens, building_density, plus neighborhood random effect.",
    7.2, 1.8, 5.0, 3.2, font_size=16
)


# Slide 4: count results
slide = prs.slides.add_slide(layout_object)
set_text(find_placeholder(slide, 0), "Primary model results: neighboring-home count", size=24, bold=True)
add_image(slide, IMG_PRIMARY_LA, 0.6, 1.5, 6.0)
add_image(slide, IMG_PRIMARY_NYC, 6.8, 1.5, 6.0)
add_textbox(slide, "LA: significant nonlinear CSI smooth\nedf ≈ 3.00, p < 0.001", 0.8, 6.55, 5.8, 0.5, font_size=15, align=PP_ALIGN.CENTER)
add_textbox(slide, "NYC: significant, close-to-linear CSI smooth\nedf ≈ 1.01, p < 0.001", 7.0, 6.55, 5.5, 0.5, font_size=15, align=PP_ALIGN.CENTER)


# Slide 5: share results
slide = prs.slides.add_slide(layout_object)
set_text(find_placeholder(slide, 0), "Sensitivity results: neighboring-home share", size=24, bold=True)
add_image(slide, IMG_SHARE_LA, 0.6, 1.5, 6.0)
add_image(slide, IMG_SHARE_NYC, 6.8, 1.5, 6.0)
add_textbox(slide, "LA: significant nonlinear CSI pattern\nedf ≈ 3.04, p < 0.001", 0.8, 6.55, 5.8, 0.5, font_size=15, align=PP_ALIGN.CENTER)
add_textbox(slide, "NYC: significant, nearly linear CSI pattern\nedf ≈ 1.00, p < 0.001", 7.0, 6.55, 5.5, 0.5, font_size=15, align=PP_ALIGN.CENTER)


# Slide 6: take-home
slide = prs.slides.add_slide(layout_content)
set_text(find_placeholder(slide, 0), "Take-home points", size=24, bold=True)
add_bullets_to_placeholder(
    find_placeholder(slide, 1),
    [
        "The neighboring-home indicator is a realized local-use/accessibility measure.",
        "It complements NDVI and distance to green space but is not direct park visitation.",
        "The metric is built from polygon-based neighboring CBG sets and aggregated to tracts.",
        "In both LA and NYC, CSI is significantly associated with neighboring-home outcomes.",
        "LA shows more nonlinear CSI patterns; NYC is closer to linear in the full-year models."
    ],
    size=20
)

prs.save(str(PPTX_PATH))
print(PPTX_PATH)
